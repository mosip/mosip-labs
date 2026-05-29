"""
Generate MOSIP Nexus documentation:
  1. Developer Guide  → MosipNexus/docs/MOSIP_Nexus_Developer_Guide.docx
  2. Business PPT     → MosipNexus/docs/MOSIP_Nexus_Presentation.pptx
"""

from pathlib import Path
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPTPt
from pptx.dml.color import RGBColor as PPTRGBColor
from pptx.enum.text import PP_ALIGN

DOCS_DIR = Path(__file__).parent

BLUE  = RGBColor(0x1a, 0x56, 0xdb)
DARK  = RGBColor(0x1a, 0x1a, 0x2e)
GRAY  = RGBColor(0x6b, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

PPT_BLUE  = PPTRGBColor(0x1a, 0x56, 0xdb)
PPT_DARK  = PPTRGBColor(0x0f, 0x17, 0x2a)
PPT_GREEN = PPTRGBColor(0x05, 0x96, 0x69)
PPT_WHITE = PPTRGBColor(0xFF, 0xFF, 0xFF)
PPT_GRAY  = PPTRGBColor(0x6b, 0x72, 0x80)
PPT_INDIGO = PPTRGBColor(0xa5, 0xb4, 0xfc)
PPT_MUTED = PPTRGBColor(0xd1, 0xd5, 0xdb)


# ─────────────────────────────────────────────────────────────────────────────
# WORD HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def set_cell_bg(cell, hex_color: str):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), hex_color)
    shd.set(qn("w:val"), "clear")
    tcPr.append(shd)


def add_heading(doc, text, level=1):
    p = doc.add_heading(text, level=level)
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run(text)
    if level == 1:
        run.font.color.rgb = BLUE
        run.font.size = Pt(18)
    elif level == 2:
        run.font.color.rgb = DARK
        run.font.size = Pt(14)
    else:
        run.font.color.rgb = GRAY
        run.font.size = Pt(12)
    return p


def add_body(doc, text):
    p = doc.add_paragraph(text)
    for run in p.runs:
        run.font.size = Pt(11)
    return p


def add_bullet(doc, text, level=0):
    p = doc.add_paragraph(text, style="List Bullet")
    p.paragraph_format.left_indent = Inches(0.25 * (level + 1))
    return p


def add_numbered(doc, text):
    return doc.add_paragraph(text, style="List Number")


def add_code(doc, text):
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Inches(0.4)
    run = p.add_run(text)
    run.font.name = "Courier New"
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0xd6, 0x33, 0x84)
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), "F3F4F6")
    shd.set(qn("w:val"), "clear")
    p._p.get_or_add_pPr().append(shd)
    return p


def add_info_box(doc, text):
    table = doc.add_table(rows=1, cols=1)
    table.style = "Table Grid"
    cell = table.cell(0, 0)
    set_cell_bg(cell, "EFF6FF")
    p = cell.paragraphs[0]
    run = p.add_run(text)
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x1e, 0x40, 0xaf)
    doc.add_paragraph()


def add_table(doc, rows, header=None):
    table = doc.add_table(rows=1 + len(rows), cols=2)
    table.style = "Light List Accent 1"
    if header:
        hdr = table.rows[0].cells
        hdr[0].text = header[0]
        hdr[1].text = header[1]
        for cell in hdr:
            set_cell_bg(cell, "1A56DB")
            for run in cell.paragraphs[0].runs:
                run.font.color.rgb = WHITE
                run.font.bold = True
    for i, (c1, c2) in enumerate(rows):
        row = table.rows[i + 1].cells
        row[0].text = c1
        row[1].text = c2
        if i % 2 == 0:
            for cell in row:
                set_cell_bg(cell, "F0F4FF")
    doc.add_paragraph()


# ─────────────────────────────────────────────────────────────────────────────
# DEVELOPER GUIDE
# ─────────────────────────────────────────────────────────────────────────────

def build_developer_guide():
    doc = Document()
    doc.styles["Normal"].font.name = "Calibri"
    doc.styles["Normal"].font.size = Pt(11)

    # Cover
    doc.add_paragraph()
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run("MOSIP Nexus")
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = BLUE

    s = doc.add_paragraph()
    s.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = s.add_run("Developer Guide")
    r2.font.size = Pt(20)
    r2.font.color.rgb = DARK

    tl = doc.add_paragraph()
    tl.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r3 = tl.add_run("AI-Powered Knowledge Assistant for MOSIP")
    r3.font.size = Pt(13)
    r3.font.color.rgb = GRAY
    r3.font.italic = True

    doc.add_paragraph()
    m = doc.add_paragraph()
    m.alignment = WD_ALIGN_PARAGRAPH.CENTER
    m.add_run("Version 1.0  |  May 2026  |  For Freshers with Zero GenAI Background")
    doc.add_page_break()

    # TOC
    add_heading(doc, "Table of Contents", 1)
    for item in [
        "1. What Is This Project?",
        "2. GenAI Fundamentals — The Concepts Behind This App",
        "3. Technology Stack",
        "4. Architecture Overview",
        "5. Project Structure",
        "6. Setting Up Your Development Environment",
        "7. Running the Application for the First Time",
        "8. How the Pipeline Works — Component by Component",
        "9. REST API Reference",
        "10. Code Walkthrough — Key Functions Explained",
        "11. Key Configuration Reference",
        "12. Making Changes",
        "13. Incremental Knowledge Updates",
        "14. Troubleshooting",
        "15. Glossary",
    ]:
        add_bullet(doc, item)
    doc.add_page_break()

    # 1. What Is This Project?
    add_heading(doc, "1. What Is This Project?", 1)
    add_heading(doc, "1.1 Background — The MOSIP Problem", 2)
    add_body(doc,
        "MOSIP (Modular Open Source Identity Platform) is a large, complex platform used by "
        "governments worldwide to build national digital identity systems. Its knowledge is "
        "spread across six different places:")
    for s in ["Official documentation at docs.mosip.io (449 pages)",
              "Community forum at community.mosip.io (thousands of Q&A threads)",
              "GitHub Issues across 86 MOSIP repositories",
              "Source code — millions of lines of Java, YAML, and properties files",
              "Confluence (internal engineering wikis)",
              "Jira (project management and bug tracking)"]:
        add_bullet(doc, s)
    add_body(doc,
        "When a developer has a question — like 'What does error IDA-MLC-009 mean?' or "
        "'Which Java class handles biometric authentication?' — they would have to manually "
        "search all of these places. This wastes hours every day.")

    add_heading(doc, "1.2 What MOSIP Nexus Does", 2)
    add_body(doc,
        "MOSIP Nexus is an AI chat assistant that lets you ask a question in plain English "
        "(or any language) and instantly get a precise answer sourced from all six knowledge "
        "bases simultaneously. Think of it as a MOSIP expert who has read everything, is "
        "available 24/7, and never gets tired.")
    add_info_box(doc,
        "Key idea: You ask one question → the system searches 70,000+ knowledge chunks across "
        "6 sources → picks the most relevant pieces → feeds them to an AI model → returns a "
        "synthesised answer with source links.")

    add_heading(doc, "1.3 Why 'Nexus'?", 2)
    add_body(doc,
        "'Nexus' means a central connection point. This app is the single hub that connects "
        "you to all MOSIP knowledge at once, so you never have to search six places manually again.")

    add_heading(doc, "1.4 Who Built This?", 2)
    add_body(doc,
        "This project was created as a personal initiative to explore how Generative AI can "
        "solve real MOSIP support and onboarding challenges. It is not an official MOSIP product "
        "— it is a working proof-of-concept built from scratch.")
    doc.add_page_break()

    # 2. GenAI Fundamentals
    add_heading(doc, "2. GenAI Fundamentals — The Concepts Behind This App", 1)
    add_body(doc,
        "If you have never worked with Generative AI, this section explains every concept you "
        "need. Read carefully — everything else builds on this.")

    add_heading(doc, "2.1 What Is a Large Language Model (LLM)?", 2)
    add_body(doc,
        "An LLM is an AI model trained on massive amounts of text from the internet, books, "
        "and code. It learns patterns in language so well that it can write, summarise, "
        "translate, and answer questions like a human expert.")
    add_body(doc,
        "MOSIP Nexus uses Llama 3.3 (70 billion parameters), provided via Groq — a fast, "
        "free-to-use AI API. You send it a question plus some context, and it writes the answer.")
    add_info_box(doc,
        "Analogy: Think of an LLM as an extremely well-read intern who has memorised millions "
        "of documents. You give it a question and some relevant pages from a manual, and it "
        "writes a clear answer for you.")

    add_heading(doc, "2.2 The Problem with LLMs Alone", 2)
    add_body(doc,
        "LLMs have a training cutoff — they don't know about your specific MOSIP version, "
        "your organisation's Confluence pages, or new GitHub issues raised last week. "
        "If you ask a raw LLM a MOSIP question, it either makes something up (hallucination) "
        "or says it doesn't know.")

    add_heading(doc, "2.3 What Is RAG (Retrieval Augmented Generation)?", 2)
    add_body(doc,
        "RAG is the technique that makes MOSIP Nexus work. Instead of asking the LLM to answer "
        "from its training memory, we:")
    add_numbered(doc, "Convert all MOSIP knowledge into a searchable index.")
    add_numbered(doc, "When you ask a question, search the index for the most relevant pieces.")
    add_numbered(doc, "Feed those pieces to the LLM as context along with your question.")
    add_numbered(doc, "The LLM writes the answer ONLY based on the provided context.")
    add_body(doc,
        "This eliminates hallucinations because the LLM is told: use only what I give you. "
        "If you don't see the answer in the context, say so.")
    add_info_box(doc,
        "Analogy: Imagine you are taking an open-book exam. RAG is the process of finding the "
        "right pages in the textbook and handing them to the student before they write the "
        "answer. The student (LLM) writes the answer based on those pages, not from memory.")

    add_heading(doc, "2.4 What Are Embeddings?", 2)
    add_body(doc,
        "An embedding is a list of numbers (a vector) that represents the meaning of a piece "
        "of text. Similar meanings produce similar vectors. For example, 'biometric "
        "authentication' and 'fingerprint login' would have vectors that are very close to "
        "each other in mathematical space.")
    add_body(doc,
        "MOSIP Nexus uses the model intfloat/multilingual-e5-base from HuggingFace. "
        "It produces 768-number vectors and supports 100+ languages — so Hindi, Tamil, and "
        "Arabic questions can also find relevant English documentation.")
    add_info_box(doc,
        "Analogy: Think of embeddings as GPS coordinates for meaning. Two sentences about the "
        "same concept have coordinates that are geographically close, so when you search for "
        "one, the other is found nearby.")

    add_heading(doc, "2.5 What Is a Vector Database?", 2)
    add_body(doc,
        "A vector database stores all those 768-number vectors and lets you find the most "
        "similar ones to a query vector very quickly — even when there are millions of vectors. "
        "MOSIP Nexus uses ChromaDB — a lightweight vector database that runs entirely on your "
        "local machine with no server required.")

    add_heading(doc, "2.6 What Is MMR (Maximal Marginal Relevance)?", 2)
    add_body(doc,
        "When searching for the top 8 relevant chunks, a naive approach just picks the 8 most "
        "similar — but they often all say the same thing. MMR is smarter: it picks chunks that "
        "are both relevant AND diverse from each other, giving the LLM richer context.")

    add_heading(doc, "2.7 Confidence Scoring", 2)
    add_body(doc,
        "Cosine distance measures how different two vectors are. Distance 0 = identical, "
        "Distance 1 = completely unrelated. MOSIP Nexus uses this to label every answer:")
    add_table(doc, [
        ("Distance < 0.25", "High confidence (similarity > 75%)"),
        ("Distance 0.25 to 0.45", "Medium confidence"),
        ("Distance > 0.45", "Low confidence — answer may be unreliable"),
    ], header=("Cosine Distance", "Confidence Label"))

    add_heading(doc, "2.8 What Is LangChain?", 2)
    add_body(doc,
        "LangChain is a Python framework that provides standard building blocks like prompt "
        "templates, retrievers, LLM connectors, and memory — so you don't have to build "
        "everything from scratch. MOSIP Nexus uses LangChain's LCEL (LangChain Expression "
        "Language) to wire together the RAG pipeline as a clean chain: prompt | LLM | parser.")
    doc.add_page_break()

    # 3. Technology Stack
    add_heading(doc, "3. Technology Stack", 1)
    add_table(doc, [
        ("LangChain (LCEL)", "Python framework for building LLM pipelines. Wires together prompt templates, retrievers, and LLMs."),
        ("Groq API", "Free cloud API that runs Llama 3.3 at very high speed. Get a free key at console.groq.com."),
        ("Llama 3.3 (70B)", "The LLM that reads context and writes answers. Developed by Meta AI."),
        ("HuggingFace Embeddings", "intfloat/multilingual-e5-base converts text to 768-dim vectors. Runs locally on CPU."),
        ("ChromaDB", "Local vector database. Stores all embeddings on disk (SQLite). No server needed."),
        ("Streamlit", "Python library for building web UIs with almost no HTML/CSS/JS."),
        ("BeautifulSoup", "HTML parser used by the docs crawler to extract text from MOSIP documentation pages."),
        ("DuckDuckGo Search", "Fallback when no MOSIP-specific answer is found — searches the public web."),
        ("langdetect", "Detects the language of the user's message so the LLM replies in the same language."),
        ("uv", "Fast Python package manager. All commands use 'uv run' instead of 'python'."),
        ("LangSmith (optional)", "Traces every LLM call with latency and token usage. Free tier at smith.langchain.com."),
    ], header=("Technology", "What It Does"))
    doc.add_page_break()

    # 4. Architecture Overview
    add_heading(doc, "4. Architecture Overview", 1)
    add_body(doc,
        "MOSIP Nexus has two distinct phases: an offline data pipeline (run once) and an online "
        "query pipeline (runs every time a user asks a question).")

    add_heading(doc, "4.1 Offline Pipeline — Building the Knowledge Base", 2)
    add_table(doc, [
        ("Stage 1: Crawl", "6 crawlers fetch content from all sources and save as JSON files in data/."),
        ("Stage 2: Chunk", "Each document is split into ~900-character overlapping text chunks."),
        ("Stage 3: Embed", "Each chunk is converted to a 768-dim vector using multilingual-e5-base."),
        ("Stage 4: Store", "Vectors and text are saved into 6 ChromaDB collections in chroma_db/."),
    ], header=("Stage", "Description"))

    add_heading(doc, "4.2 Online Pipeline — Answering a Question", 2)
    add_table(doc, [
        ("1. Input", "User types a question in the Streamlit chat UI."),
        ("2. Language detection", "langdetect identifies the language. Replies in that language for the session."),
        ("3. Greeting check", "Pure greetings get a friendly reply without searching the knowledge base."),
        ("4. Question condensing", "If there is chat history, the LLM rewrites the follow-up as a standalone query."),
        ("5. MMR retrieval", "The query is embedded and searched across all 6 ChromaDB collections using MMR."),
        ("6. Confidence scoring", "Best cosine distance across all collections becomes the confidence label."),
        ("7. Answer generation", "Chunks + question + history are sent to Groq Llama 3.3. Answer is returned."),
        ("8. Fallback", "If the question is unrelated to MOSIP, DuckDuckGo web search is tried."),
        ("9. Response", "Answer + sources + confidence badge + related threads are displayed in the UI."),
    ], header=("Step", "Description"))
    doc.add_page_break()

    # 5. Project Structure
    add_heading(doc, "5. Project Structure", 1)
    add_body(doc, "All code lives inside the MosipNexus/ folder:")
    add_table(doc, [
        ("config/settings.py", "Central configuration. All constants, thresholds, API keys, and paths. Import this everywhere."),
        ("crawler/docs_crawler.py", "Crawls all 449 pages of docs.mosip.io using the sitemap. Outputs data/mosip_docs.json."),
        ("crawler/community_crawler.py", "Fetches Q&A threads from community.mosip.io via Discourse JSON API."),
        ("crawler/github_crawler.py", "Fetches closed issues from all 86 MOSIP GitHub repos via GitHub REST API."),
        ("crawler/code_crawler.py", "Fetches Java/YAML/properties source files via GitHub Tree API."),
        ("crawler/confluence_crawler.py", "Fetches pages from Confluence spaces (requires credentials)."),
        ("crawler/jira_crawler.py", "Fetches Jira tickets (optional, requires credentials)."),
        ("crawler/state.py", "Saves crawl state (URL hashes, topic IDs) to track what has been crawled."),
        ("ingestion/store.py", "Reads all JSON files, chunks, embeds, and upserts into ChromaDB."),
        ("retrieval/retriever.py", "MMR search across all ChromaDB collections + confidence scoring."),
        ("retrieval/dedup.py", "Detects if a question is similar to an existing community thread."),
        ("chain/query_engine.py", "Core RAG pipeline: greeting → condense → retrieve → answer → fallback."),
        ("memory/session.py", "Manages per-session chat history as LangChain message objects."),
        ("notifications/email_notifier.py", "Sends email alerts for unanswerable questions (optional SMTP config)."),
        ("app/app.py", "Streamlit chat UI with language lock, export, sidebar, confidence badges."),
        ("run_update.py", "Incremental update runner — fetches only what changed since last run."),
        ("data/", "Generated JSON files from crawlers. Not committed to git."),
        ("chroma_db/", "Generated ChromaDB vector store. Not committed to git."),
        (".env", "Your secret API keys. Never committed to git."),
    ], header=("File / Folder", "Purpose"))
    doc.add_page_break()

    # 6. Setup
    add_heading(doc, "6. Setting Up Your Development Environment", 1)
    add_heading(doc, "6.1 Prerequisites", 2)
    add_table(doc, [
        ("Python 3.13", "Required. Download from python.org."),
        ("uv", "Package manager. Install with: pip install uv"),
        ("Git", "For cloning the repository."),
        ("Groq API key", "Free. Sign up at console.groq.com and create an API key."),
        ("GitHub Token (optional)", "Enables 5,000 req/hr. Generate at github.com / Settings / Developer settings."),
        ("8 GB free disk space", "For ChromaDB vector store and embedding model downloads (~300 MB model)."),
    ], header=("Requirement", "Notes"))

    add_heading(doc, "6.2 Installation Steps", 2)
    for label, code in [
        ("Clone the repository", "git clone <repo-url>\ncd LangChainUpdated"),
        ("Install all dependencies", "uv sync"),
        ("Copy the environment file template", "copy MosipNexus\\.env.example MosipNexus\\.env"),
        ("Open .env and add your Groq API key", 'GROQ_API_KEY="<YOUR_GROQ_API_KEY>"'),
    ]:
        add_body(doc, label + ":")
        add_code(doc, code)

    add_heading(doc, "6.3 Environment Variables Reference", 2)
    add_body(doc,
        "Create MosipNexus/.env and populate the variables below. Only GROQ_API_KEY is "
        "required to run the app. All others are optional and enable additional features.")
    add_table(doc, [
        ("GROQ_API_KEY", "REQUIRED. Your Groq API key. Free at console.groq.com. Without this the app cannot generate answers."),
        ("HF_TOKEN", "Optional. HuggingFace token. Avoids rate limits when downloading the multilingual-e5-base embedding model. Get from huggingface.co/settings/tokens."),
        ("LANGCHAIN_TRACING_V2", "Optional. Set to true to enable LangSmith tracing. Traces every LLM call with latency and token usage."),
        ("LANGCHAIN_API_KEY", "Optional. LangSmith API key. Required when LANGCHAIN_TRACING_V2=true. Free at smith.langchain.com."),
        ("LANGCHAIN_PROJECT", "Optional. LangSmith project name. Defaults to mosip-nexus. All traces appear under this project."),
        ("GITHUB_TOKEN", "Optional but recommended. GitHub personal access token. Raises the API rate limit from 60 to 5,000 requests/hour for the GitHub crawler. Generate at github.com > Settings > Developer settings > Personal access tokens."),
        ("CONFLUENCE_URL", "Optional. Your Atlassian Confluence base URL, e.g. https://your-org.atlassian.net/wiki. Required for the Confluence crawler."),
        ("CONFLUENCE_USER", "Optional. Your Atlassian account email address. Used with CONFLUENCE_TOKEN for Basic Auth."),
        ("CONFLUENCE_TOKEN", "Optional. Atlassian API token. Generate at id.atlassian.com > Security > API tokens."),
        ("CONFLUENCE_SPACE_KEYS", "Optional. Comma-separated Confluence space keys to crawl, e.g. QT,ENGG,PMS."),
        ("JIRA_URL", "Optional. Your Atlassian Jira base URL, e.g. https://your-org.atlassian.net."),
        ("JIRA_USER", "Optional. Your Atlassian account email address for Jira authentication."),
        ("JIRA_TOKEN", "Optional. Atlassian API token for Jira (can be the same token as CONFLUENCE_TOKEN)."),
        ("JIRA_PROJECT_KEYS", "Optional. Comma-separated Jira project keys to crawl, e.g. MOSIP,MISP."),
        ("SMTP_HOST", "Optional. SMTP server hostname for email alerts, e.g. smtp.gmail.com."),
        ("SMTP_PORT", "Optional. SMTP port number. Typically 587 for TLS."),
        ("SMTP_USER", "Optional. Email address used to send alert emails."),
        ("SMTP_PASSWORD", "Optional. App password for the SMTP account (for Gmail: generate an App Password under your Google Account security settings)."),
        ("NOTIFY_EMAIL", "Optional. Recipient email address for low-confidence answer alerts."),
    ], header=("Variable", "Description"))
    add_info_box(doc,
        "Minimum .env to get started:\n"
        "GROQ_API_KEY=\"<YOUR_GROQ_API_KEY>\"\n\n"
        "Full example with all optional features:\n"
        "GROQ_API_KEY=\"<YOUR_GROQ_API_KEY>\"\n"
        "HF_TOKEN=\"<YOUR_HF_TOKEN>\"\n"
        "LANGCHAIN_TRACING_V2=true\n"
        "LANGCHAIN_API_KEY=\"<YOUR_LANGCHAIN_API_KEY>\"\n"
        "LANGCHAIN_PROJECT=mosip-nexus\n"
        "GITHUB_TOKEN=\"<YOUR_GITHUB_TOKEN>\"\n"
        "CONFLUENCE_URL=https://your-org.atlassian.net/wiki\n"
        "CONFLUENCE_USER=you@example.com\n"
        "CONFLUENCE_TOKEN=\"<YOUR_CONFLUENCE_TOKEN>\"\n"
        "CONFLUENCE_SPACE_KEYS=QT,ENGG,PMS")

    add_heading(doc, "6.4 Verify Setup", 2)
    add_code(doc, 'uv run python -c "from chain.query_engine import ask; print(\'Setup OK\')"')
    doc.add_page_break()

    # 7. Running
    add_heading(doc, "7. Running the Application for the First Time", 1)
    add_info_box(doc,
        "Important: Steps 1 and 2 only need to be done once (or when you want to refresh the "
        "knowledge base). Step 3 is the daily startup command.")

    add_heading(doc, "7.1 Step 1 — Run the Crawlers", 2)
    add_body(doc,
        "Run from the LangChainUpdated root directory. Each crawler is independent and can "
        "be run in separate terminals simultaneously.")
    for cmd, note in [
        ("uv run python MosipNexus/crawler/docs_crawler.py", "~5 min — 449 documentation pages"),
        ("uv run python MosipNexus/crawler/community_crawler.py", "~30 min — thousands of forum threads"),
        ("uv run python MosipNexus/crawler/github_crawler.py", "~15 min with token — 86 repos"),
        ("uv run python MosipNexus/crawler/code_crawler.py", "~60 min — source code files"),
    ]:
        add_body(doc, f"  {note}:")
        add_code(doc, cmd)

    add_heading(doc, "7.2 Step 2 — Build the Vector Index", 2)
    add_body(doc,
        "Reads JSON files, generates embeddings (CPU), stores in ChromaDB. "
        "Expect 2–6 hours for a full index. Progress shown per batch.")
    add_code(doc, "uv run python MosipNexus/ingestion/store.py")
    add_info_box(doc,
        "To rebuild from scratch (e.g. after changing the embedding model):\n"
        "Remove-Item -Recurse -Force MosipNexus/chroma_db")

    add_heading(doc, "7.3 Step 3 — Launch the App", 2)
    add_code(doc, "uv run --no-sync streamlit run MosipNexus/app/app.py")
    add_body(doc, "Open http://localhost:8501 in your browser.")
    doc.add_page_break()

    # 8. Pipeline components
    add_heading(doc, "8. How the Pipeline Works — Component by Component", 1)

    add_heading(doc, "8.1 Crawlers (crawler/)", 2)
    add_body(doc,
        "Each crawler is a standalone Python script. Output is a JSON array where each object "
        "has at minimum: url, title, and content. The docs crawler uses MOSIP's sitemap.xml to "
        "discover pages. The community crawler uses the Discourse JSON API. The GitHub crawler "
        "uses the GitHub REST API and auto-discovers all 86 MOSIP organisation repos.")

    add_heading(doc, "8.2 Ingestion (ingestion/store.py)", 2)
    add_numbered(doc, "Load — Read the JSON file.")
    add_numbered(doc, "Prepare — Convert each item to a LangChain Document with metadata (source_type, URL, title).")
    add_numbered(doc, "Chunk — Split into 900-character chunks with 150-character overlap.")
    add_numbered(doc, "Embed — Generate a 768-dim vector per chunk using multilingual-e5-base.")
    add_numbered(doc, "Store — Upsert into the matching ChromaDB collection.")
    add_info_box(doc,
        "Why chunk? LLMs have a context window limit. Chunking breaks long documents into "
        "pieces small enough to fit, while overlap ensures no information is lost at boundaries.")

    add_heading(doc, "8.3 Retriever (retrieval/retriever.py)", 2)
    add_numbered(doc, "Embed the query using the same multilingual-e5-base model.")
    add_numbered(doc, "Run MMR search on each ChromaDB collection.")
    add_numbered(doc, "Merge results from all collections.")
    add_numbered(doc, "Compute best cosine distance across all collections → confidence score.")
    add_numbered(doc, "Return the ranked list of Document objects.")
    add_body(doc,
        "Error code queries (e.g. IDA-MLC-009) and code/class queries get a 3x retrieval boost "
        "from the code collection, since those answers are most likely in error constant files "
        "or service implementations.")

    add_heading(doc, "8.4 Query Engine (chain/query_engine.py)", 2)
    add_numbered(doc, "Pure greeting check — reply directly, skip RAG.")
    add_numbered(doc, "Question condensing — rewrite follow-up questions as standalone queries using the LLM.")
    add_numbered(doc, "Retrieval — call retrieve() from retriever.py.")
    add_numbered(doc, "System prompt building — select error format, code format, or general format.")
    add_numbered(doc, "Answer generation — QA chain (prompt | LLM | parser) generates the answer.")
    add_numbered(doc, "Fallback — if [NO_DOC_SOURCE] marker detected, try DuckDuckGo web search.")
    add_numbered(doc, "Source classification — determine mosip_docs / community / github / code / mixed.")

    add_heading(doc, "8.5 Streamlit App (app/app.py)", 2)
    for item in [
        "Chat input and message rendering",
        "Language detection with auto-detect and language lock",
        "Explicit language instruction detection ('reply in Tamil' → popup confirmation)",
        "Sidebar: New Chat, Export Chat, session info, language buttons",
        "Export — self-contained HTML report with sources and confidence badges",
        "Confidence badges, source expanders, similar question expanders",
        "Auto-scroll to latest message with floating Latest message button",
        "Email notification when confidence is low",
        "Duplicate question detection surface",
    ]:
        add_bullet(doc, item)
    doc.add_page_break()

    # 9. REST API Reference
    add_heading(doc, "9. REST API Reference", 1)
    add_body(doc,
        "MOSIP Nexus exposes a production-ready REST API built with FastAPI. Every feature "
        "available in the Streamlit chat UI is also available as an HTTP endpoint — making it "
        "straightforward to integrate MOSIP Nexus into portals, Slack bots, mobile apps, or "
        "any custom internal tool.")

    add_heading(doc, "9.1 Starting the API Server", 2)
    add_body(doc, "Run from the LangChainUpdated root directory (the ChromaDB index must already be built):")
    add_code(doc, "uv run uvicorn MosipNexus.api.main:app --host 0.0.0.0 --port 8000 --reload")
    add_table(doc, [
        ("Swagger UI", "http://localhost:8000/docs — interactive browser-based API explorer, try every endpoint without writing code"),
        ("ReDoc", "http://localhost:8000/redoc — clean, readable API documentation"),
        ("API Base URL", "http://localhost:8000"),
    ], header=("Resource", "URL"))
    add_info_box(doc,
        "The Streamlit UI (port 8501) and the REST API (port 8000) can run simultaneously — "
        "they share the same ChromaDB index and embedding model.")

    add_heading(doc, "9.2 Endpoint Reference", 2)
    add_table(doc, [
        ("POST /chat", "Main RAG endpoint. Send a question, get answer + sources + confidence. Pass session_id to continue a conversation."),
        ("POST /batch", "Ask up to 10 questions in one HTTP call. All share the same session context — each answer is aware of the previous."),
        ("POST /search", "Raw semantic search across all ChromaDB collections. Returns matching chunks without calling the LLM. Useful for building custom integrations."),
        ("POST /similar", "Check if a similar community forum thread already exists for the question. No LLM is called."),
        ("POST /session", "Create a new empty session. Returns a session_id to use in subsequent /chat calls."),
        ("GET /session/{id}/history", "Retrieve the full Q&A history (all turns) for a session."),
        ("GET /sessions", "List all currently active sessions with their turn count and detected language."),
        ("DELETE /session/{id}", "Clear a session and its history. Equivalent to clicking New Chat in the UI."),
        ("POST /feedback", "Rate an answer as positive or negative for a specific turn in a session."),
        ("POST /export/{id}", "Export a full session as JSON (structured data) or HTML (styled report identical to the UI export)."),
        ("GET /stats", "Aggregated usage metrics — total queries, confidence distribution, source types, and languages used."),
        ("GET /health", "Returns ChromaDB collection chunk counts and active session count. Use for monitoring."),
    ], header=("Endpoint", "Description"))

    add_heading(doc, "9.3 Example — Ask a Question (/chat)", 2)
    add_body(doc, "Request:")
    add_code(doc,
        "POST http://localhost:8000/chat\n"
        "Content-Type: application/json\n"
        "\n"
        "{\n"
        "  \"question\": \"What does error IDA-MLC-009 mean?\",\n"
        "  \"language\": \"English\"\n"
        "}"
    )
    add_body(doc, "Response:")
    add_code(doc,
        "{\n"
        "  \"answer\": \"IDA-MLC-009 indicates a license expiry error in ID Authentication...\",\n"
        "  \"sources\": [\n"
        "    {\"source\": \"https://docs.mosip.io/...\", \"title\": \"Error Codes\", \"source_type\": \"docs\"}\n"
        "  ],\n"
        "  \"source_type\": \"mixed\",\n"
        "  \"confidence\": \"high\",\n"
        "  \"similar_questions\": [],\n"
        "  \"session_id\": \"a3f2b1c4-8e91-4d2a-b3f7-12c456789abc\"\n"
        "}"
    )
    add_info_box(doc,
        "Pass the session_id from the response back in your next /chat request to maintain "
        "conversation context. The API remembers all previous questions in the session — "
        "follow-up questions like 'What about its error handling?' work correctly.")

    add_heading(doc, "9.4 Example — Multi-turn Conversation", 2)
    add_body(doc, "First call (no session_id — a new one is created automatically):")
    add_code(doc,
        "POST /chat\n"
        "{\"question\": \"What is Key Manager in MOSIP?\", \"language\": \"English\"}\n"
        "→ returns session_id: \"abc-123\""
    )
    add_body(doc, "Follow-up call (pass the session_id back):")
    add_code(doc,
        "POST /chat\n"
        "{\"question\": \"Which modules depend on it?\", \"session_id\": \"abc-123\"}\n"
        "→ API understands 'it' = Key Manager from previous turn"
    )

    add_heading(doc, "9.5 Example — Batch Questions (/batch)", 2)
    add_body(doc, "Process multiple questions in one HTTP call:")
    add_code(doc,
        "POST http://localhost:8000/batch\n"
        "Content-Type: application/json\n"
        "\n"
        "{\n"
        "  \"questions\": [\n"
        "    \"What is MOSIP?\",\n"
        "    \"What are its main modules?\",\n"
        "    \"Which module handles biometrics?\"\n"
        "  ],\n"
        "  \"language\": \"English\"\n"
        "}"
    )
    add_body(doc,
        "Returns all three answers in one response. Each answer is aware of the previous — "
        "the third question can reference context from the first two.")

    add_heading(doc, "9.6 Example — Usage Statistics (/stats)", 2)
    add_body(doc, "GET http://localhost:8000/stats — useful for stakeholder reporting:")
    add_code(doc,
        "{\n"
        "  \"total_queries\": 142,\n"
        "  \"active_sessions\": 5,\n"
        "  \"positive_feedback\": 38,\n"
        "  \"negative_feedback\": 4,\n"
        "  \"avg_sources_per_query\": 3.2,\n"
        "  \"confidence_distribution\": {\"high\": 89, \"medium\": 41, \"low\": 12},\n"
        "  \"source_type_distribution\": {\"mixed\": 95, \"mosip_docs\": 30, \"community\": 17},\n"
        "  \"language_distribution\": {\"English\": 120, \"Tamil\": 14, \"Hindi\": 8}\n"
        "}"
    )
    doc.add_page_break()

    # 10. Code Walkthrough
    add_heading(doc, "10. Code Walkthrough — Key Functions Explained", 1)
    add_body(doc,
        "This chapter walks through the most important functions in the codebase with actual "
        "code snippets and line-by-line explanations. Read this after understanding the architecture "
        "in Chapter 8 — it shows exactly HOW each concept is implemented in Python.")

    # ── 9.1 memory/session.py ─────────────────────────────────────────────────
    add_heading(doc, "10.1  memory/session.py — SessionMemory Class", 2)
    add_body(doc,
        "This is the simplest file in the project. It stores one user's conversation history "
        "for a single browser session. Streamlit keeps one instance of this in st.session_state.")
    add_code(doc,
        "@dataclass\n"
        "class SessionMemory:\n"
        "    messages: list[BaseMessage] = field(default_factory=list)\n"
        "    language: str = \"English\"\n"
        "    lang_code: str = \"en\"\n"
        "\n"
        "    def add_turn(self, question: str, answer: str) -> None:\n"
        "        self.messages.append(HumanMessage(content=question))\n"
        "        self.messages.append(AIMessage(content=answer))\n"
        "\n"
        "    def clear(self) -> None:\n"
        "        self.messages.clear()\n"
        "        self.language = \"English\"\n"
        "        self.lang_code = \"en\""
    )
    add_table(doc, [
        ("@dataclass", "Python decorator — auto-generates __init__, __repr__ from the class fields. No need to write a constructor manually."),
        ("messages: list[BaseMessage]", "Stores the full conversation. Each turn adds two items: HumanMessage (question) and AIMessage (answer). LangChain understands these types natively."),
        ("field(default_factory=list)", "Ensures each instance gets its own empty list, not a shared one. A common Python gotcha with mutable defaults."),
        ("language / lang_code", "Tracks the currently detected language. 'English' / 'en' is the default. Updated by langdetect on every query."),
        ("add_turn()", "Called after every Q&A exchange. Appends both messages to history so the condenser chain has context for follow-up questions."),
        ("clear()", "Called when the user clicks 'New Chat'. Resets everything to initial state."),
    ], header=("Code Element", "Explanation"))

    # ── 9.2 retrieval/dedup.py ────────────────────────────────────────────────
    add_heading(doc, "10.2  retrieval/dedup.py — find_similar_question()", 2)
    add_body(doc,
        "Before running the full RAG pipeline, the app checks whether the question was already "
        "asked in the community forum. If yes, it surfaces the existing thread instead of generating "
        "a new answer from scratch.")
    add_code(doc,
        "def find_similar_question(query: str) -> dict | None:\n"
        "    embeddings = get_embeddings()           # reuse shared embedding model\n"
        "    store = Chroma(\n"
        "        persist_directory=CHROMA_DIR,\n"
        "        embedding_function=embeddings,\n"
        "        collection_name=COMMUNITY_COLLECTION,\n"
        "    )\n"
        "    results = store.similarity_search_with_score(\n"
        "        query,\n"
        "        k=1,\n"
        "        filter={\"post_type\": \"question\"},   # search questions only, not answers\n"
        "    )\n"
        "    doc, distance = results[0]\n"
        "    similarity = 1.0 - distance            # convert distance → similarity\n"
        "    if similarity < DEDUP_THRESHOLD:       # DEDUP_THRESHOLD = 0.88\n"
        "        return None\n"
        "    return {\"title\": ..., \"source\": ..., \"similarity_score\": round(similarity, 3)}"
    )
    add_table(doc, [
        ("similarity_search_with_score()", "Like similarity_search() but also returns the cosine distance for each result. We need the distance to compute the similarity score."),
        ("filter={\"post_type\": \"question\"}", "ChromaDB metadata filter — only matches chunks that were tagged as questions during ingestion. This avoids matching answer text, which would give false positives."),
        ("similarity = 1.0 - distance", "ChromaDB returns cosine DISTANCE (0=identical, 1=unrelated). We want SIMILARITY (1=identical, 0=unrelated), so we flip it."),
        ("DEDUP_THRESHOLD = 0.88", "If similarity is above 88%, we consider it the same question. Set in config/settings.py. Tunable."),
        ("Return None", "None means 'no duplicate found — proceed with normal RAG'. The app checks for None and skips the dedup surface if not found."),
    ], header=("Code Element", "Explanation"))

    # ── 9.3 retrieval/retriever.py ────────────────────────────────────────────
    add_heading(doc, "10.3  retrieval/retriever.py — retrieve()", 2)
    add_body(doc,
        "This is the search engine of MOSIP Nexus. It queries all ChromaDB collections and "
        "returns the most relevant knowledge chunks for a given query, along with a confidence label.")
    add_code(doc,
        "def retrieve(query: str, k: int = RETRIEVAL_K) -> tuple[list[Document], str]:\n"
        "    docs_store, community_store = _get_stores()\n"
        "\n"
        "    # MMR search on each collection\n"
        "    docs_retriever = docs_store.as_retriever(\n"
        "        search_type=\"mmr\",\n"
        "        search_kwargs={\"k\": k, \"fetch_k\": RETRIEVAL_FETCH_K},\n"
        "    )\n"
        "    doc_results       = docs_retriever.invoke(query)\n"
        "    community_results = community_retriever.invoke(query)\n"
        "\n"
        "    # Code collection: 3x boost for error code / class questions\n"
        "    _is_targeted = _ERROR_CODE_RE.search(query) or _CODE_QUERY_RE.search(query)\n"
        "    code_k = CODE_RETRIEVAL_K * 3 if _is_targeted else CODE_RETRIEVAL_K\n"
        "\n"
        "    # Confidence: best cosine distance across all collections\n"
        "    best_distance = float(\"inf\")\n"
        "    for store, _ in _score_stores:\n"
        "        top = store.similarity_search_with_score(query, k=1)\n"
        "        if top:\n"
        "            best_distance = min(best_distance, top[0][1])\n"
        "\n"
        "    if best_distance < CONFIDENCE_HIGH:    # 0.25\n"
        "        confidence = \"high\"\n"
        "    elif best_distance < CONFIDENCE_MEDIUM: # 0.45\n"
        "        confidence = \"medium\"\n"
        "    else:\n"
        "        confidence = \"low\"\n"
        "\n"
        "    return doc_results + community_results + github_results + code_results, confidence"
    )
    add_table(doc, [
        ("search_type=\"mmr\"", "Maximal Marginal Relevance — returns diverse, relevant chunks instead of 8 near-identical ones."),
        ("fetch_k=30", "MMR first fetches 30 candidates, then picks the best k=8 from them. Larger fetch_k = better diversity."),
        ("as_retriever().invoke(query)", "LangChain abstracts the embedding + search into one call. Internally: embed query → vector search → return Document objects."),
        ("_ERROR_CODE_RE.search(query)", "Regex that detects MOSIP error codes like IDA-MLC-009. If found, the code collection gets 3x more chunks since error constants live in Java files."),
        ("best_distance = min(...)", "Takes the closest match across ALL collections. Not just docs — if a GitHub issue has the perfect answer, confidence should be high."),
        ("float('inf') as starting value", "A safe starting point for a minimum search. Any real distance will be smaller than infinity, so the first result always wins."),
        ("Return tuple", "Returns both the documents (for LLM context) and the confidence string (for the UI badge). Two things in one return."),
    ], header=("Code Element", "Explanation"))

    # ── 9.4 chain/query_engine.py — _build() ─────────────────────────────────
    add_heading(doc, "10.4  chain/query_engine.py — _build() and the Condenser Chain", 2)
    add_body(doc,
        "_build() sets up two objects that are reused for every query: the LLM instance and the "
        "condenser chain. The condenser rewrites follow-up questions into standalone search queries.")
    add_code(doc,
        "def _build() -> None:\n"
        "    global _llm, _condenser\n"
        "    _llm = ChatGroq(model=GROQ_MODEL, api_key=GROQ_API_KEY)\n"
        "\n"
        "    condense_prompt = ChatPromptTemplate.from_messages([\n"
        "        (\"system\",\n"
        "         \"Given the conversation history and the latest question, rewrite the \"\n"
        "         \"question as a standalone search query. Return ONLY the rewritten question.\"),\n"
        "        MessagesPlaceholder(\"chat_history\"),\n"
        "        (\"human\", \"{input}\"),\n"
        "    ])\n"
        "    _condenser = condense_prompt | _llm | StrOutputParser()"
    )
    add_table(doc, [
        ("global _llm, _condenser", "Module-level variables. Python's way of sharing state across function calls without a class. _build() is called only once (lazy init)."),
        ("ChatGroq(...)", "LangChain wrapper around the Groq API. Behaves like any other LangChain LLM — you can swap it for ChatOpenAI with no other changes."),
        ("ChatPromptTemplate.from_messages()", "Builds a prompt with multiple message roles: system (instructions), chat_history (prior conversation), human (current question)."),
        ("MessagesPlaceholder(\"chat_history\")", "A slot in the prompt that gets filled with the actual chat history list at call time. The LLM sees all prior turns."),
        ("condense_prompt | _llm | StrOutputParser()", "This is LCEL — LangChain Expression Language. The | operator chains steps: prompt → LLM → parse output as string. No boilerplate needed."),
        ("StrOutputParser()", "Extracts just the text string from the LLM's response object. Without it, you'd get an AIMessage object instead of a plain string."),
    ], header=("Code Element", "Explanation"))

    # ── 9.5 chain/query_engine.py — _build_system_prompt() ───────────────────
    add_heading(doc, "10.5  chain/query_engine.py — _build_system_prompt()", 2)
    add_body(doc,
        "This function builds the system prompt that tells the LLM how to behave. "
        "It dynamically adjusts the answer format based on the type of question.")
    add_code(doc,
        "def _build_system_prompt(\n"
        "    language: str,\n"
        "    has_greeting: bool,\n"
        "    is_error_query: bool = False,\n"
        "    is_code_query: bool = False,\n"
        ") -> str:\n"
        "    if is_error_query:\n"
        "        answer_format = (\n"
        "            \"For this error code question, structure your answer as:\\n\"\n"
        "            \"  1. What the error means\\n\"\n"
        "            \"  2. Most likely root cause(s)\\n\"\n"
        "            \"  3. Step-by-step resolution\\n\"\n"
        "            \"  4. How to verify the fix worked\\n\"\n"
        "        )\n"
        "    elif is_code_query:\n"
        "        answer_format = (\n"
        "            \"Name the exact class(es) from the retrieved context directly.\\n\"\n"
        "            \"State what the class does and which package it belongs to.\\n\"\n"
        "        )\n"
        "    else:\n"
        "        answer_format = \"Match format to question type — definitions get explanations...\"\n"
        "\n"
        "    return (\n"
        "        f\"You are a senior MOSIP engineer...\\n\\n\"\n"
        "        f\"RULE #1 — FABRICATION IS FORBIDDEN:\\n\"\n"
        "        f\"Every class name, method name, CLI command you write MUST appear \"\n"
        "        f\"word-for-word in the retrieved context below...\\n\\n\"\n"
        "        f\"Retrieved context:\\n{{context}}\"\n"
        "    )"
    )
    add_table(doc, [
        ("is_error_query", "True when the question contains a MOSIP error code pattern (IDA-MLC-009). Detected via regex: r'\\b[A-Z]{2,}-[A-Z]{2,}-\\d{3,}\\b'. Gets a structured 4-part answer format."),
        ("is_code_query", "True when the question asks about a class, method, file, or configuration. Detected via regex. Gets a format that names the class first and immediately."),
        ("answer_format variable", "Injected into the system prompt. This means the LLM gets a different instruction depending on the question type — without any if/else in the LLM call itself."),
        ("RULE #1 — FABRICATION IS FORBIDDEN", "The most important rule in the prompt. Placed at the very top so the LLM reads it first. Explicitly bans inventing class names, commands, or repo names not in the context."),
        ("{context}", "A placeholder in the prompt string filled in at call time with the retrieved document chunks. The LLM is shown the relevant MOSIP knowledge right in the prompt."),
        ("f-string with {{context}}", "Double braces {{ }} are how you write a literal { } inside an f-string. Without doubling, Python would try to evaluate {context} as a variable."),
    ], header=("Code Element", "Explanation"))

    # ── 9.6 chain/query_engine.py — ask() ────────────────────────────────────
    add_heading(doc, "10.6  chain/query_engine.py — ask() — The Full RAG Pipeline", 2)
    add_body(doc,
        "ask() is the central function of the entire project. Every user message goes through it. "
        "Understanding this function means understanding the entire RAG pipeline.")
    add_code(doc,
        "def ask(question: str, chat_history: list, language: str = \"English\") -> dict:\n"
        "\n"
        "    # Step 1: Pure greeting — skip RAG entirely\n"
        "    if _is_greeting(question):\n"
        "        response = llm.invoke([SystemMessage(...), HumanMessage(content=question)])\n"
        "        return {\"answer\": response.content, \"source_type\": \"chat\", ...}\n"
        "\n"
        "    # Step 2: Condense follow-up questions into standalone queries\n"
        "    standalone = (\n"
        "        _condenser.invoke({\"input\": question, \"chat_history\": chat_history})\n"
        "        if chat_history else question\n"
        "    )\n"
        "\n"
        "    # Step 3: Retrieve from all ChromaDB collections\n"
        "    docs, confidence = retrieve(standalone)\n"
        "    context = \"\\n\\n\".join(d.page_content for d in docs)\n"
        "\n"
        "    # Step 4: Build and run the QA chain\n"
        "    system_prompt = _build_system_prompt(language, has_greeting, is_error_query, is_code_query)\n"
        "    qa_prompt = ChatPromptTemplate.from_messages([\n"
        "        (\"system\", system_prompt),\n"
        "        MessagesPlaceholder(\"chat_history\"),\n"
        "        (\"human\", \"{input}\"),\n"
        "    ])\n"
        "    qa_chain = qa_prompt | llm | StrOutputParser()\n"
        "    answer = qa_chain.invoke({\n"
        "        \"input\": question,\n"
        "        \"chat_history\": chat_history,\n"
        "        \"context\": context,         # <-- retrieved chunks injected here\n"
        "    })\n"
        "\n"
        "    # Step 5: Fallback — web search if LLM says [NO_DOC_SOURCE]\n"
        "    if _NO_SOURCE_MARKER in answer:\n"
        "        web_hits = _web_search(question)\n"
        "        if web_hits:\n"
        "            return {\"answer\": web_response.content, \"source_type\": \"web\", ...}\n"
        "\n"
        "    return {\"answer\": answer, \"sources\": sources, \"confidence\": confidence, ...}"
    )
    add_table(doc, [
        ("_is_greeting(question)", "Checks if the message is a pure greeting ('hi', 'thanks', 'bye'). If yes, the LLM replies conversationally without touching the knowledge base — saves tokens and avoids irrelevant retrieval."),
        ("_condenser.invoke(...) if chat_history else question", "The condenser rewrites 'What about its error handling?' into 'How does MOSIP Registration Processor handle errors?' Only run when there IS prior history — no point condensing a fresh question."),
        ("context = '\\n\\n'.join(d.page_content for d in docs)", "Concatenates all retrieved document chunks into one big string that is injected into the system prompt. The LLM reads this as its knowledge source."),
        ("qa_chain = qa_prompt | llm | StrOutputParser()", "LCEL chain assembled fresh for each call. prompt fills in {context}/{input}/{chat_history} → LLM generates → parser extracts string."),
        ("qa_chain.invoke({...})", "Single call that runs the entire chain. LangChain handles sending the filled prompt to Groq and returning the response."),
        ("[NO_DOC_SOURCE] marker", "If the question is completely unrelated to MOSIP (e.g. 'What is the weather?'), the system prompt instructs the LLM to output this marker. The code detects it and triggers web search."),
        ("Return dict", "Always returns a dict with answer, sources, source_type, confidence, similar_questions. The Streamlit UI unpacks this to render badges, expanders, and source links."),
    ], header=("Code Element", "Explanation"))

    # ── 9.7 ingestion/store.py — ingest() ────────────────────────────────────
    add_heading(doc, "10.7  ingestion/store.py — ingest()", 2)
    add_body(doc,
        "This function takes a list of LangChain Documents, chunks them, and stores them "
        "in a ChromaDB collection. It is the core of the build-once knowledge pipeline.")
    add_code(doc,
        "def ingest(documents: list[Document], collection_name: str,\n"
        "           embeddings: HuggingFaceEmbeddings) -> None:\n"
        "\n"
        "    splitter = RecursiveCharacterTextSplitter(\n"
        "        chunk_size=CHUNK_SIZE,      # 900 characters\n"
        "        chunk_overlap=CHUNK_OVERLAP # 150 characters\n"
        "    )\n"
        "    chunks = splitter.split_documents(documents)\n"
        "    # chunks is now a bigger list — each Document may produce multiple chunks\n"
        "\n"
        "    vectorstore = None\n"
        "    for i in range(total_batches):          # process 100 chunks at a time\n"
        "        batch = chunks[i * 100 : (i+1) * 100]\n"
        "        if vectorstore is None:\n"
        "            vectorstore = Chroma.from_documents(  # creates collection\n"
        "                documents=batch,\n"
        "                embedding=embeddings,\n"
        "                persist_directory=CHROMA_DIR,\n"
        "                collection_name=collection_name,\n"
        "            )\n"
        "        else:\n"
        "            vectorstore.add_documents(batch)      # appends to existing"
    )
    add_table(doc, [
        ("RecursiveCharacterTextSplitter", "Splits text into chunks trying to break at paragraph boundaries first, then sentences, then words, then characters. 'Recursive' means it tries larger separators first."),
        ("chunk_size=900 / chunk_overlap=150", "Each chunk is at most 900 characters. The last 150 characters of one chunk are repeated at the start of the next. This ensures context at boundaries is not lost."),
        ("splitter.split_documents()", "Takes a list of Documents and returns a bigger list of Documents. Metadata (source URL, title, etc.) is copied to all child chunks automatically."),
        ("Chroma.from_documents() on first batch", "Creates the collection and inserts the first 100 chunks. Internally: embed all 100 → store vectors + text in SQLite."),
        ("vectorstore.add_documents() on subsequent batches", "Appends to the already-created collection. This is why the first call is different from subsequent ones."),
        ("Batching (100 chunks at a time)", "Embedding 38,000 chunks at once would exhaust RAM. Batching keeps memory usage constant regardless of collection size."),
    ], header=("Code Element", "Explanation"))

    # ── 9.8 app.py — Session state and query flow ─────────────────────────────
    add_heading(doc, "10.8  app/app.py — Session State and Query Flow", 2)
    add_body(doc,
        "Streamlit re-runs the entire app.py script from top to bottom on every user interaction. "
        "st.session_state is how you persist data across those re-runs.")
    add_code(doc,
        "# Session state is initialised once per browser session\n"
        "if \"memory\" not in st.session_state:\n"
        "    st.session_state.memory = SessionMemory()    # chat history\n"
        "if \"messages\" not in st.session_state:\n"
        "    st.session_state.messages = []               # display history\n"
        "if \"lang_locked\" not in st.session_state:\n"
        "    st.session_state.lang_locked = False         # language lock toggle\n"
        "if \"pending_lang\" not in st.session_state:\n"
        "    st.session_state.pending_lang = None         # awaiting user confirmation\n"
        "\n"
        "# Chat input is captured BEFORE rendering (so sidebar shows updated language)\n"
        "query = st.chat_input(\"Ask anything about MOSIP...\")\n"
        "if query:\n"
        "    lang_instruction = _detect_lang_instruction(query)  # 'reply in Tamil'?\n"
        "    if lang_instruction and lang_instruction[0] != memory.lang_code:\n"
        "        st.session_state.pending_lang = lang_instruction  # show popup\n"
        "    elif not st.session_state.lang_locked:\n"
        "        lang_result = _detect_language(query, memory.lang_code)\n"
        "        if lang_result:\n"
        "            memory.set_language(*lang_result)"
    )
    add_table(doc, [
        ("if 'memory' not in st.session_state", "The guard pattern for Streamlit. Without this check, session state would be reset to a new SessionMemory() on every re-run (every keypress)."),
        ("st.session_state.memory = SessionMemory()", "Stores the SessionMemory object. Because it's in session_state, it survives page re-runs and retains all chat history."),
        ("messages vs memory.messages", "Two separate lists. memory.messages = LangChain message objects passed to the LLM chains. st.session_state.messages = display dicts for rendering in the UI with sources and confidence."),
        ("query = st.chat_input()", "Streamlit's chat input widget. Returns None when nothing is typed, or the string when submitted. The app processes the query in the same script run."),
        ("lang_instruction detection", "Checks if the user explicitly said 'reply in Tamil' or similar. If yes, sets pending_lang for popup confirmation instead of auto-switching silently."),
        ("elif not lang_locked", "Auto-detect only runs when the language is NOT locked. Once the user locks to Tamil, even English questions get Tamil answers."),
        ("*lang_result (unpacking)", "lang_result is a tuple (lang_code, lang_name). The * unpacks it into two positional arguments for set_language(lang_code, lang_name)."),
    ], header=("Code Element", "Explanation"))

    add_code(doc,
        "# After rendering messages — language popup confirmation\n"
        "if st.session_state.pending_lang:\n"
        "    _plc, _pln = st.session_state.pending_lang\n"
        "    st.info(f\"Stay in **{_pln}** for the rest of this conversation?\")\n"
        "    col1, col2 = st.columns(2)\n"
        "    with col1:\n"
        "        if st.button(f\"Yes, lock to {_pln}\", key=\"_lang_yes\"):\n"
        "            memory.set_language(_plc, _pln)\n"
        "            st.session_state.lang_locked = True\n"
        "            st.session_state.pending_lang = None\n"
        "            st.rerun()     # re-render with updated sidebar\n"
        "    with col2:\n"
        "        if st.button(\"No, just this reply\", key=\"_lang_no\"):\n"
        "            st.session_state.pending_lang = None\n"
        "            st.rerun()"
    )
    add_table(doc, [
        ("st.session_state.pending_lang", "Holds a (lang_code, lang_name) tuple while waiting for user confirmation. Set when a language instruction is detected. Cleared when user clicks Yes or No."),
        ("st.columns(2)", "Creates two side-by-side columns. Each column renders independently. Used here to place two buttons next to each other."),
        ("key='_lang_yes'", "Every interactive widget in Streamlit needs a unique key. Without it, Streamlit cannot distinguish between two buttons with the same label across re-runs."),
        ("st.rerun()", "Forces Streamlit to re-run the entire script immediately. Used after state changes so the UI reflects the new state (e.g., updated sidebar language label)."),
    ], header=("Code Element", "Explanation"))

    add_code(doc,
        "# After generating an answer — update state and trigger sidebar refresh\n"
        "if query:\n"
        "    with st.chat_message(\"assistant\"):\n"
        "        with st.spinner(\"Searching docs and community forum...\"):\n"
        "            result = ask(query, memory.messages, memory.language)\n"
        "\n"
        "        answer      = result[\"answer\"]\n"
        "        confidence  = result[\"confidence\"]\n"
        "        source_type = result[\"source_type\"]\n"
        "        st.markdown(answer)\n"
        "\n"
        "    st.session_state.messages.append({\n"
        "        \"role\": \"assistant\",\n"
        "        \"content\": answer,\n"
        "        \"sources\": sources,\n"
        "        \"confidence\": confidence,\n"
        "    })\n"
        "    memory.add_turn(query, answer)  # update LangChain history\n"
        "    st.rerun()                      # so export button sees the new message"
    )
    add_table(doc, [
        ("with st.chat_message('assistant')", "Context manager that renders everything inside it inside a chat bubble with the assistant avatar. Works for 'user' and 'assistant' roles."),
        ("with st.spinner(...)", "Shows a loading spinner while the indented code runs. Removed automatically when the block exits. The ask() call is slow (LLM API call), so the spinner gives visual feedback."),
        ("ask(query, memory.messages, memory.language)", "The three arguments: the raw question, the full LangChain history (for condensing), and the current session language (for the system prompt)."),
        ("st.session_state.messages.append({...})", "Saves the answer with its metadata to the DISPLAY history. This is what gets rendered on the next re-run and included in the HTML export."),
        ("memory.add_turn(query, answer)", "Saves to the LangChain history. This is what gets passed to the condenser on the next question so it can handle follow-ups."),
        ("st.rerun() at the end", "Forces a re-render so the sidebar 'Export Chat' button appears (it only shows when messages exist) and the export captures the new message."),
    ], header=("Code Element", "Explanation"))

    doc.add_page_break()

    # 11. Configuration
    add_heading(doc, "11. Key Configuration Reference", 1)
    add_table(doc, [
        ("RETRIEVAL_K = 8", "Top-K chunks returned per collection per query."),
        ("RETRIEVAL_FETCH_K = 30", "MMR candidate pool size. Larger = more diverse results."),
        ("CHUNK_SIZE = 900", "Max characters per chunk."),
        ("CHUNK_OVERLAP = 150", "Overlap between chunks. Prevents losing context at boundaries."),
        ("CONFIDENCE_HIGH = 0.25", "Cosine distance threshold for High confidence badge."),
        ("CONFIDENCE_MEDIUM = 0.45", "Cosine distance threshold for Medium confidence badge."),
        ("DEDUP_THRESHOLD = 0.88", "Cosine similarity above which a question is a duplicate."),
        ("EMBED_MODEL", "HuggingFace model ID. Changing this requires a full re-embed."),
        ("GROQ_MODEL", "Groq LLM model ID. Currently llama-3.3-70b-versatile."),
        ("CODE_RETRIEVAL_K = 6", "Chunks from code collection. Boosted 3x for error code or class queries."),
    ], header=("Setting", "Description"))
    doc.add_page_break()

    # 12. Making Changes
    add_heading(doc, "12. Making Changes", 1)

    add_heading(doc, "12.1 Changing the Answer Style / Prompt", 2)
    add_body(doc,
        "Edit _build_system_prompt() in chain/query_engine.py. The prompt has three sections: "
        "persona, answer format (changes per query type), and strict rules. "
        "After saving, Streamlit auto-reloads — no restart needed.")

    add_heading(doc, "12.2 Adding a New Knowledge Source", 2)
    add_numbered(doc, "Write a new crawler that outputs a JSON file with url, title, and content fields.")
    add_numbered(doc, "Add a collection name to config/settings.py (e.g. NEW_COLLECTION = 'mosip_new').")
    add_numbered(doc, "Add a prepare_new_documents() function in ingestion/store.py.")
    add_numbered(doc, "Add _try_load_optional_store(NEW_COLLECTION, embeddings) in retrieval/retriever.py.")
    add_numbered(doc, "Add the store to the retrieval loop in retrieve().")

    add_heading(doc, "12.3 Changing the LLM", 2)
    add_body(doc,
        "Replace ChatGroq in chain/query_engine.py with any LangChain-compatible LLM class "
        "(e.g. ChatOpenAI, ChatAnthropic). The rest of the pipeline is unchanged.")

    add_heading(doc, "12.4 Changing the Embedding Model", 2)
    add_body(doc,
        "Change EMBED_MODEL in config/settings.py. Delete chroma_db/ and re-run "
        "ingestion/store.py. The pipeline version hash updates automatically.")
    doc.add_page_break()

    # 13. Incremental Updates
    add_heading(doc, "13. Incremental Knowledge Updates", 1)
    add_body(doc,
        "After the initial crawl, run_update.py fetches only new or changed content "
        "by comparing against data/crawl_state.json.")
    add_code(doc, "uv run python MosipNexus/run_update.py")
    add_table(doc, [
        ("Docs", "Re-fetches pages whose content hash has changed."),
        ("Community", "Fetches topics with ID greater than the last seen max."),
        ("GitHub", "Fetches issues with number greater than the last seen max per repo."),
        ("Code", "Re-fetches files that have changed based on GitHub commit SHAs."),
    ], header=("Source", "Update Strategy"))
    doc.add_page_break()

    # 14. Troubleshooting
    add_heading(doc, "14. Troubleshooting", 1)
    add_table(doc, [
        ("App crashes: GROQ_API_KEY not set", "Create MosipNexus/.env with GROQ_API_KEY='your_key'."),
        ("'Access is denied' error on startup", "Run with: uv run --no-sync streamlit run MosipNexus/app/app.py"),
        ("'Expecting value: line 1 column 1'", "A data JSON file is empty. Re-run the corresponding crawler."),
        ("Embedding model download hangs", "Set HF_TOKEN in .env or check your internet connection."),
        ("Language not updating for short queries", "Short non-ASCII queries (<5 chars) need longer text to trigger detection."),
        ("Groq 429 rate limit error", "Free tier is 100K tokens/day. Wait a few minutes. The UI shows a warning."),
        ("ChromaDB dimension mismatch", "You changed EMBED_MODEL. Delete chroma_db/ and re-run ingestion/store.py."),
        ("Community crawler stops early", "Rate limiting from community.mosip.io. Increase CRAWL_DELAY_SECS in settings.py."),
    ], header=("Error", "Fix"))
    doc.add_page_break()

    # 15. Glossary
    add_heading(doc, "15. Glossary", 1)
    add_table(doc, [
        ("LLM", "Large Language Model — an AI trained on massive text corpora to understand and generate language."),
        ("RAG", "Retrieval Augmented Generation — retrieve relevant docs first, then ask the LLM to answer."),
        ("Embedding", "A vector (list of numbers) representing the semantic meaning of a piece of text."),
        ("Vector Database", "A database optimised for storing and searching embeddings by similarity."),
        ("ChromaDB", "The vector database used in this project — runs locally with no server."),
        ("MMR", "Maximal Marginal Relevance — retrieval algorithm balancing relevance and diversity."),
        ("Cosine Distance", "Measure of difference between two vectors. 0 = identical, 1 = unrelated."),
        ("Chunk", "A small segment of a larger document created by splitting for embedding."),
        ("LCEL", "LangChain Expression Language — syntax for composing chains with the | operator."),
        ("Groq", "Cloud API service running LLMs at very high speed. Free tier available."),
        ("Llama 3.3", "The LLM used — 70 billion parameter model by Meta AI."),
        ("Streamlit", "Python library for building web apps without HTML/CSS/JS."),
        ("Discourse", "Forum software running community.mosip.io with a public JSON API."),
        ("Confidence Score", "High/Medium/Low label derived from cosine distance of the best retrieved chunk."),
        ("Hallucination", "When an LLM invents facts not present in its training data or context."),
        ("uv", "Fast Python package manager used in this project. Replaces pip."),
    ], header=("Term", "Definition"))

    out = DOCS_DIR / "MOSIP_Nexus_Developer_Guide.docx"
    doc.save(str(out))
    print(f"Developer Guide saved: {out}")


# ─────────────────────────────────────────────────────────────────────────────
# PPT HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def add_slide(prs, layout_index=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_index])

def set_bg(slide, r, g, b):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = PPTRGBColor(r, g, b)

def txbox(slide, text, left, top, width, height,
          size=18, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    color = color or PPT_WHITE
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PPTPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def rect(slide, left, top, width, height, fill, line=None):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def blist(slide, items, left, top, width, height, size=13, color=None, char="▸"):
    color = color or PPT_MUTED
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = PPTPt(3)
        run = p.add_run()
        run.text = f"{char}  {item}"
        run.font.size = PPTPt(size)
        run.font.color.rgb = color


# ─────────────────────────────────────────────────────────────────────────────
# PRESENTATION
# ─────────────────────────────────────────────────────────────────────────────

def build_ppt():
    prs = Presentation()
    W = prs.slide_width  = PInches(13.33)
    H = prs.slide_height = PInches(7.5)

    # Light professional palette
    L_BG    = (0xF5, 0xF7, 0xFB)                  # slide background
    L_NAVY  = PPTRGBColor(0x1E, 0x3A, 0x8A)       # header bars, main accent
    L_BLUE  = PPTRGBColor(0x1D, 0x4E, 0xD8)       # card titles, secondary accent
    L_CARD  = PPTRGBColor(0xEB, 0xF5, 0xFF)       # card background
    L_CBRD  = PPTRGBColor(0xBC, 0xD4, 0xF5)       # card border
    L_DARK  = PPTRGBColor(0x1F, 0x29, 0x37)       # body text
    L_GRAY  = PPTRGBColor(0x6B, 0x72, 0x80)       # secondary / footer text
    L_WHITE = PPTRGBColor(0xFF, 0xFF, 0xFF)        # white (on dark bars)
    L_GREEN = PPTRGBColor(0x06, 0x5F, 0x46)       # dark green
    L_GBG   = PPTRGBColor(0xEC, 0xFD, 0xF5)       # light green card bg
    L_GBRD  = PPTRGBColor(0x6E, 0xE7, 0xB7)       # green card border
    L_ORANG = PPTRGBColor(0xB4, 0x5F, 0x09)       # orange (future roadmap)
    L_OBG   = PPTRGBColor(0xFF, 0xF7, 0xED)       # light orange bg
    L_OBRD  = PPTRGBColor(0xFD, 0xBA, 0x74)       # orange border

    # ── Slide 1: Title ──
    s = add_slide(prs)
    set_bg(s, *L_BG)
    rect(s, 0, 0, W, PInches(0.14), L_NAVY)
    rect(s, 0, PInches(2.82), W, PInches(0.06), L_BLUE)
    rect(s, 0, PInches(7.1), W, PInches(0.4), L_NAVY)
    txbox(s, "MOSIP Nexus", PInches(1), PInches(0.5), PInches(11.3), PInches(1.6),
          size=54, bold=True, color=L_NAVY, align=PP_ALIGN.CENTER)
    txbox(s, "AI-Powered Knowledge Assistant for MOSIP",
          PInches(1), PInches(2.92), PInches(11.3), PInches(0.8),
          size=22, color=L_BLUE, align=PP_ALIGN.CENTER, italic=True)
    txbox(s, "A Personal Initiative  |  May 2026",
          PInches(1), PInches(3.8), PInches(11.3), PInches(0.55),
          size=14, color=L_GRAY, align=PP_ALIGN.CENTER)
    txbox(s, "One AI assistant. Six knowledge sources. Any language. Instant answers.",
          PInches(1.5), PInches(5.3), PInches(10.3), PInches(0.65),
          size=16, color=L_GREEN, align=PP_ALIGN.CENTER, italic=True)
    txbox(s, "MOSIP Nexus  |  Confidential",
          PInches(1), PInches(7.13), PInches(11.3), PInches(0.28),
          size=10, color=L_WHITE, align=PP_ALIGN.CENTER)

    # ── Slide 2: The Challenge ──
    s = add_slide(prs)
    set_bg(s, *L_BG)
    rect(s, 0, 0, W, PInches(1.25), L_NAVY)
    txbox(s, "The Challenge", PInches(0.5), PInches(0.18), PInches(12), PInches(0.72),
          size=30, bold=True, color=L_WHITE)
    txbox(s, "MOSIP's knowledge is scattered across 6 different places. Finding one answer can take hours.",
          PInches(0.5), PInches(0.88), PInches(12), PInches(0.4),
          size=13, color=PPTRGBColor(0xBC, 0xD4, 0xF5), italic=True)
    sources = [
        ("docs.mosip.io", "449 pages of\nofficial docs"),
        ("community.mosip.io", "Thousands of\nQ&A threads"),
        ("GitHub Issues", "86 repos with\nclosed issues"),
        ("Source Code", "Millions of lines\nJava / YAML"),
        ("Confluence", "Internal\nengineering wikis"),
        ("Jira", "Project tickets\n& bug reports"),
    ]
    bw = PInches(1.9)
    gap = PInches(0.18)
    sx = PInches(0.4)
    for i, (title, desc) in enumerate(sources):
        x = sx + i * (bw + gap)
        rect(s, x, PInches(2.3), bw, PInches(2.6), L_CARD, L_CBRD)
        txbox(s, title, x + PInches(0.1), PInches(2.45), bw - PInches(0.2), PInches(0.55),
              size=11, bold=True, color=L_BLUE, align=PP_ALIGN.CENTER)
        txbox(s, desc, x + PInches(0.1), PInches(3.05), bw - PInches(0.2), PInches(1.5),
              size=11, color=L_DARK, align=PP_ALIGN.CENTER)
    txbox(s,
          "A support engineer searching for one error code answer wastes 2-4 hours per incident.",
          PInches(0.5), PInches(5.2), PInches(12), PInches(0.55),
          size=13, color=L_ORANG, italic=True)

    # ── Slide 3: The Solution ──
    s = add_slide(prs)
    set_bg(s, *L_BG)
    rect(s, 0, 0, W, PInches(1.25), L_GREEN)
    txbox(s, "The Solution: MOSIP Nexus", PInches(0.5), PInches(0.18), PInches(12), PInches(0.72),
          size=30, bold=True, color=L_WHITE)
    txbox(s,
          "One chat assistant that searches all six sources simultaneously and returns "
          "a precise, cited answer in seconds — in any language.",
          PInches(0.5), PInches(0.88), PInches(12), PInches(0.4),
          size=13, color=PPTRGBColor(0xD1, 0xFA, 0xE5), italic=True)
    cols = [
        ("Ask Anything", "Type your MOSIP question in plain language. No need to know which source has the answer."),
        ("Instant Retrieval", "AI searches 70,000+ knowledge chunks from all 6 sources at once using semantic search."),
        ("Cited Answer", "Get a precise answer with links to the exact doc page, GitHub issue, or community thread."),
        ("Any Language", "Ask in Tamil, Hindi, French, Arabic — answer comes back in the same language automatically."),
    ]
    cw = PInches(3.0)
    for i, (title, desc) in enumerate(cols):
        x = PInches(0.4) + i * PInches(3.2)
        rect(s, x, PInches(2.3), cw, PInches(3.6), L_GBG, L_GBRD)
        txbox(s, title, x + PInches(0.1), PInches(2.45), cw - PInches(0.2), PInches(0.55),
              size=14, bold=True, color=L_GREEN)
        txbox(s, desc, x + PInches(0.1), PInches(3.05), cw - PInches(0.2), PInches(2.6),
              size=12, color=L_DARK)

    # ── Slide 4: How It Works ──
    s = add_slide(prs)
    set_bg(s, *L_BG)
    rect(s, 0, 0, W, PInches(1.25), L_NAVY)
    txbox(s, "How It Works — In 4 Simple Steps", PInches(0.5), PInches(0.18), PInches(12), PInches(0.8),
          size=30, bold=True, color=L_WHITE)
    steps = [
        ("1", "You Ask", "Type any MOSIP question in natural language — in any language you prefer."),
        ("2", "We Search", "AI searches 70,000+ chunks across 6 sources using semantic (meaning-based) search."),
        ("3", "AI Answers", "Groq Llama 3.3 reads the top results and writes a precise, grounded answer."),
        ("4", "You Get", "Answer + confidence score + source links + related threads — in under 10 seconds."),
    ]
    bw2 = PInches(2.9)
    for i, (num, title, desc) in enumerate(steps):
        x = PInches(0.4) + i * PInches(3.2)
        rect(s, x, PInches(1.5), bw2, PInches(4.85), L_CARD, L_CBRD)
        rect(s, x + PInches(0.9), PInches(1.65), PInches(1.1), PInches(1.0), L_NAVY)
        txbox(s, num, x + PInches(0.9), PInches(1.68), PInches(1.1), PInches(0.85),
              size=32, bold=True, color=L_WHITE, align=PP_ALIGN.CENTER)
        txbox(s, title, x + PInches(0.1), PInches(2.8), bw2 - PInches(0.2), PInches(0.5),
              size=15, bold=True, color=L_BLUE)
        txbox(s, desc, x + PInches(0.1), PInches(3.35), bw2 - PInches(0.2), PInches(2.7),
              size=12, color=L_DARK)
        if i < 3:
            txbox(s, "→", x + bw2 + PInches(0.05), PInches(2.9), PInches(0.4), PInches(0.5),
                  size=24, bold=True, color=L_NAVY, align=PP_ALIGN.CENTER)

    # ── Slide 5: Knowledge Sources ──
    s = add_slide(prs)
    set_bg(s, *L_BG)
    rect(s, 0, 0, W, PInches(1.25), L_NAVY)
    txbox(s, "Six Knowledge Sources — Searched Simultaneously", PInches(0.5), PInches(0.18),
          PInches(12), PInches(0.85), size=28, bold=True, color=L_WHITE)
    srcs = [
        ("MOSIP Docs", "6,094 chunks", "449 pages from docs.mosip.io/1.2.0 — full product documentation"),
        ("Community Forum", "11,236 chunks", "Accepted answers and high-voted replies from community.mosip.io"),
        ("GitHub Issues", "2,513 chunks", "Closed issues and resolutions from all 86 MOSIP repositories"),
        ("Source Code", "~38,000 chunks", "Java, YAML, properties files — error constants, service implementations"),
        ("Confluence", "12,264 chunks", "Internal engineering wikis (QT, ENGG, PMS spaces)"),
        ("Jira", "Optional", "Project tickets and bug reports (configurable per organisation)"),
    ]
    for i, (title, count, desc) in enumerate(srcs):
        row = i // 2
        col = i % 2
        x = PInches(0.4) + col * PInches(6.5)
        y = PInches(1.38) + row * PInches(1.95)
        rect(s, x, y, PInches(6.2), PInches(1.75), L_CARD, L_CBRD)
        txbox(s, f"{title}  ({count})", x + PInches(0.15), y + PInches(0.15),
              PInches(5.8), PInches(0.45), size=14, bold=True, color=L_BLUE)
        txbox(s, desc, x + PInches(0.15), y + PInches(0.65),
              PInches(5.8), PInches(0.95), size=12, color=L_DARK)

    # ── Slide 6: Key Features ──
    s = add_slide(prs)
    set_bg(s, *L_BG)
    rect(s, 0, 0, W, PInches(1.25), L_NAVY)
    txbox(s, "Key Features", PInches(0.5), PInches(0.18), PInches(12), PInches(0.85),
          size=30, bold=True, color=L_WHITE)
    features = [
        ("Multilingual", "Ask in Tamil, Hindi, Arabic, French — replies in the same language. Auto-detects. Lock to preferred language."),
        ("Chat Memory", "Follow-up questions work naturally. Full conversation context retained."),
        ("Source Attribution", "Every answer cites the exact doc page, forum thread, or GitHub issue."),
        ("Confidence Score", "High / Medium / Low badge shows how well the knowledge base matched your question."),
        ("Duplicate Detection", "Surfaces existing community threads before generating a new answer."),
        ("Web Fallback", "If MOSIP knowledge base has no answer, searches the public web and discloses the source."),
        ("Expert Escalation", "Low-confidence answers auto-notify the team and offer a Contact Expert button."),
        ("Export Report", "Download the full chat as a professional HTML report with sources and badges."),
        ("Incremental Updates", "Knowledge base updates fetch only what changed — no full rebuild needed."),
        ("Observability", "LangSmith tracing shows every prompt and LLM call with latency (optional)."),
    ]
    for i, (title, desc) in enumerate(features):
        row = i // 2
        col = i % 2
        x = PInches(0.4) + col * PInches(6.5)
        y = PInches(1.32) + row * PInches(1.18)
        txbox(s, f"  {title}", x, y, PInches(6.2), PInches(0.4),
              size=14, bold=True, color=L_BLUE)
        txbox(s, desc, x + PInches(0.1), y + PInches(0.38), PInches(6.0), PInches(0.68),
              size=11, color=L_DARK)

    # ── Slide 7: Multilingual ──
    s = add_slide(prs)
    set_bg(s, *L_BG)
    rect(s, 0, 0, W, PInches(1.25), L_NAVY)
    txbox(s, "Truly Multilingual", PInches(0.5), PInches(0.18), PInches(12), PInches(0.72),
          size=30, bold=True, color=L_WHITE)
    txbox(s, "MOSIP serves countries across the world. MOSIP Nexus speaks to users in their language.",
          PInches(0.5), PInches(0.88), PInches(12), PInches(0.4),
          size=13, color=PPTRGBColor(0xBC, 0xD4, 0xF5), italic=True)
    langs = [
        ("Tamil", "MOSIP என்றால் என்ன?"),
        ("Hindi", "MOSIP kya hai?"),
        ("French", "Qu'est-ce que le MOSIP?"),
        ("Arabic", "What is MOSIP? (in Arabic)"),
        ("German", "Was ist MOSIP?"),
        ("Portuguese", "O que e o MOSIP?"),
    ]
    lw = PInches(3.9)
    for i, (lang, q) in enumerate(langs):
        row = i // 3
        col = i % 3
        x = PInches(0.5) + col * PInches(4.2)
        y = PInches(1.9) + row * PInches(1.95)
        rect(s, x, y, lw, PInches(1.75), L_CARD, L_CBRD)
        txbox(s, lang, x + PInches(0.15), y + PInches(0.12), lw - PInches(0.3), PInches(0.45),
              size=15, bold=True, color=L_BLUE)
        txbox(s, f'"{q}"', x + PInches(0.15), y + PInches(0.62), lw - PInches(0.3), PInches(0.85),
              size=12, color=L_DARK, italic=True)
    txbox(s,
          "Auto-detects language with 95% confidence  |  Lock language for full session  |  100+ languages supported",
          PInches(0.5), PInches(6.55), PInches(12), PInches(0.5),
          size=12, color=L_GRAY, align=PP_ALIGN.CENTER)

    # ── Slide 8: Technical Architecture ──
    s = add_slide(prs)
    set_bg(s, *L_BG)
    rect(s, 0, 0, W, PInches(1.25), L_NAVY)
    txbox(s, "Technical Architecture (Simplified)", PInches(0.5), PInches(0.18),
          PInches(12), PInches(0.85), size=28, bold=True, color=L_WHITE)
    rect(s, PInches(0.3), PInches(1.38), PInches(5.5), PInches(5.6), L_GBG, L_GBRD)
    txbox(s, "OFFLINE — Build Knowledge Base (Once)", PInches(0.4), PInches(1.43),
          PInches(5.2), PInches(0.5), size=12, bold=True, color=L_GREEN)
    blist(s, [
        "6 Crawlers → docs, forum, GitHub, code, Confluence, Jira",
        "Text Splitter → 900-char chunks with 150-char overlap",
        "Embedding Model (multilingual-e5-base) → 768-dim vectors",
        "ChromaDB → stores 70,000+ vectors locally on disk",
    ], PInches(0.4), PInches(1.98), PInches(5.2), PInches(4.5), size=12, color=L_DARK)
    rect(s, PInches(6.2), PInches(1.38), PInches(6.8), PInches(5.6), L_CARD, L_CBRD)
    txbox(s, "ONLINE — Answer Every Question", PInches(6.3), PInches(1.43),
          PInches(6.5), PInches(0.5), size=12, bold=True, color=L_BLUE)
    blist(s, [
        "User asks → language detected automatically",
        "Question condensed with chat history (LLM)",
        "MMR search across all 6 ChromaDB collections",
        "Top-K chunks → context for the LLM",
        "Groq Llama 3.3 generates answer from context only",
        "DuckDuckGo fallback if question is not in knowledge base",
        "Answer + sources + confidence badge → Streamlit UI",
    ], PInches(6.3), PInches(1.98), PInches(6.6), PInches(4.5), size=12, color=L_DARK)
    txbox(s, "=>", PInches(5.75), PInches(3.85), PInches(0.45), PInches(0.6),
          size=28, bold=True, color=L_NAVY, align=PP_ALIGN.CENTER)
    rect(s, 0, PInches(6.92), W, PInches(0.58), L_NAVY)
    txbox(s, "LangChain  .  Groq Llama 3.3  .  ChromaDB  .  HuggingFace  .  Streamlit  .  Python 3.13",
          PInches(0.5), PInches(6.97), PInches(12), PInches(0.4),
          size=11, color=L_WHITE, align=PP_ALIGN.CENTER)

    # ── Slide 9: REST API ──
    s = add_slide(prs)
    set_bg(s, *L_BG)
    rect(s, 0, 0, W, PInches(1.25), L_NAVY)
    txbox(s, "REST API — Ready for Integration Today", PInches(0.5), PInches(0.18),
          PInches(12), PInches(0.72), size=28, bold=True, color=L_WHITE)
    txbox(s, "Every feature in the chat UI is also available as an HTTP endpoint with Swagger UI built-in.",
          PInches(0.5), PInches(0.88), PInches(12), PInches(0.4),
          size=13, color=PPTRGBColor(0xBC, 0xD4, 0xF5), italic=True)

    # Left panel — endpoints table
    rect(s, PInches(0.3), PInches(1.38), PInches(8.6), PInches(5.7), L_CARD, L_CBRD)
    txbox(s, "Endpoints", PInches(0.45), PInches(1.43), PInches(8.2), PInches(0.4),
          size=12, bold=True, color=L_BLUE)
    endpoints = [
        ("POST  /chat",         "Ask any MOSIP question — returns answer, sources, confidence + session_id"),
        ("POST  /batch",        "Ask up to 10 questions in one call, all sharing session context"),
        ("POST  /search",       "Raw semantic search across ChromaDB — no LLM called"),
        ("POST  /similar",      "Find matching community threads for a question"),
        ("GET   /stats",        "Usage metrics — queries, confidence levels, source types, languages"),
        ("GET   /health",       "ChromaDB chunk counts and active session count"),
        ("POST  /feedback",     "Rate answers positive / negative per session turn"),
        ("POST  /export/{id}",  "Export full session as JSON or HTML report"),
    ]
    for i, (ep, desc) in enumerate(endpoints):
        y = PInches(1.95) + i * PInches(0.58)
        txbox(s, ep, PInches(0.48), y, PInches(2.4), PInches(0.52),
              size=10, bold=True, color=L_NAVY)
        txbox(s, desc, PInches(2.95), y, PInches(5.8), PInches(0.52),
              size=10, color=L_DARK)

    # Right panel — Swagger + use cases
    rect(s, PInches(9.2), PInches(1.38), PInches(3.8), PInches(2.75), L_GBG, L_GBRD)
    txbox(s, "Integration Use Cases", PInches(9.32), PInches(1.43), PInches(3.5), PInches(0.4),
          size=12, bold=True, color=L_GREEN)
    blist(s, [
        "MOSIP Admin Portal widget",
        "Slack / Teams bot",
        "Mobile app integration",
        "Custom internal dashboards",
    ], PInches(9.32), PInches(1.92), PInches(3.5), PInches(1.9), size=12, color=L_DARK)

    rect(s, PInches(9.2), PInches(4.25), PInches(3.8), PInches(2.83), L_CARD, L_CBRD)
    txbox(s, "Swagger UI — built in", PInches(9.32), PInches(4.3), PInches(3.5), PInches(0.4),
          size=12, bold=True, color=L_BLUE)
    txbox(s, "http://localhost:8000/docs\n\nTest every endpoint from\nyour browser — no code needed.",
          PInches(9.32), PInches(4.78), PInches(3.5), PInches(1.2),
          size=11, color=L_DARK)
    txbox(s, "Start with one command:", PInches(9.32), PInches(5.98), PInches(3.5), PInches(0.3),
          size=9, bold=True, color=L_GRAY)
    txbox(s, "uv run uvicorn MosipNexus.api.main:app\n--port 8000 --reload",
          PInches(9.32), PInches(6.3), PInches(3.5), PInches(0.55),
          size=9, color=L_NAVY)

    # ── Slide 10: Business Value ──
    s = add_slide(prs)
    set_bg(s, *L_BG)
    rect(s, 0, 0, W, PInches(1.25), L_GREEN)
    txbox(s, "Business Value & Impact", PInches(0.5), PInches(0.18), PInches(12), PInches(0.72),
          size=30, bold=True, color=L_WHITE)
    metrics = [
        ("Hours -> Seconds", "An answer that took 2-4 hours of manual search now takes under 10 seconds."),
        ("Barrier-Free Access", "Non-English speaking teams in member countries get support in their native language."),
        ("Self-Service Support", "Developers and implementers resolve issues without escalating every question."),
        ("Knowledge Preserved", "Community accepted answers and GitHub resolutions are surfaced, not buried."),
        ("Intelligent Escalation", "Low-confidence answers alert the team automatically — nothing falls through."),
        ("Always Up-to-Date", "Incremental crawlers keep the knowledge base current without a full rebuild."),
    ]
    for i, (title, desc) in enumerate(metrics):
        row = i // 2
        col = i % 2
        x = PInches(0.4) + col * PInches(6.5)
        y = PInches(1.42) + row * PInches(1.9)
        rect(s, x, y, PInches(6.2), PInches(1.75), L_GBG, L_GBRD)
        txbox(s, f"  {title}", x, y + PInches(0.12), PInches(6.2), PInches(0.45),
              size=15, bold=True, color=L_GREEN)
        txbox(s, desc, x + PInches(0.15), y + PInches(0.62), PInches(5.8), PInches(0.95),
              size=12, color=L_DARK)

    # ── Slide 10: Roadmap ──
    s = add_slide(prs)
    set_bg(s, *L_BG)
    rect(s, 0, 0, W, PInches(1.25), L_NAVY)
    txbox(s, "Roadmap — What's Next", PInches(0.5), PInches(0.18), PInches(12), PInches(0.85),
          size=30, bold=True, color=L_WHITE)
    phases = [
        ("Phase 1 — Done",   L_GREEN,  L_GBG,  L_GBRD, [
            "RAG with 6 knowledge sources",
            "Multilingual chat with language lock",
            "Confidence scoring & source attribution",
            "Export report & duplicate detection",
            "Expert escalation & email notifications",
            "Incremental knowledge updates",
        ]),
        ("Phase 2 — Next",   L_BLUE,   L_CARD, L_CBRD, [
            "REST API for MOSIP portal integration",
            "Slack / Teams bot",
            "Voice input support",
            "Auto-scheduled weekly knowledge refresh",
            "User feedback per answer (thumbs up/down)",
        ]),
        ("Phase 3 — Future", L_ORANG,  L_OBG,  L_OBRD, [
            "Analytics: top questions, gap analysis",
            "Fine-tuned MOSIP embedding model",
            "Multi-tenant for different organisations",
            "MOSIP Admin Portal integration",
        ]),
    ]
    for i, (title, hdr, bg, brd, items) in enumerate(phases):
        x = PInches(0.3) + i * PInches(4.35)
        rect(s, x, PInches(1.38), PInches(4.1), PInches(5.75), bg, brd)
        rect(s, x, PInches(1.38), PInches(4.1), PInches(0.55), hdr)
        txbox(s, title, x + PInches(0.1), PInches(1.41), PInches(3.9), PInches(0.45),
              size=13, bold=True, color=L_WHITE)
        blist(s, items, x + PInches(0.15), PInches(2.03), PInches(3.85), PInches(4.8),
              size=12, color=L_DARK)

    # ── Slide 11: About This Initiative ──
    s = add_slide(prs)
    set_bg(s, *L_BG)
    rect(s, 0, 0, W, PInches(1.25), L_NAVY)
    txbox(s, "About This Initiative", PInches(0.5), PInches(0.18), PInches(12), PInches(0.72),
          size=30, bold=True, color=L_WHITE)
    txbox(s,
          "This project was built independently as a proof-of-concept to explore how "
          "Generative AI can solve a real, day-to-day challenge for the MOSIP ecosystem.",
          PInches(0.5), PInches(1.32), PInches(12), PInches(0.85),
          size=15, color=L_BLUE)
    blist(s, [
        "Built from scratch — no external AI product was licensed or purchased.",
        "Uses only open-source models and free-tier APIs (Groq, HuggingFace).",
        "Fully functional today — not a prototype, not a mockup.",
        "Knowledge base contains 70,000+ chunks from real MOSIP sources.",
        "Answering real MOSIP questions with cited sources and confidence scoring.",
        "Ready for pilot deployment with any MOSIP team or partner organisation.",
    ], PInches(0.5), PInches(2.28), PInches(12), PInches(3.5), size=14, color=L_DARK)
    rect(s, PInches(0.5), PInches(6.12), PInches(12.3), PInches(0.92), L_GBG, L_GBRD)
    txbox(s,
          '"I built this because I believe AI can genuinely reduce the friction for MOSIP adopters worldwide."',
          PInches(0.6), PInches(6.26), PInches(12.0), PInches(0.6),
          size=13, color=L_GREEN, italic=True, align=PP_ALIGN.CENTER)

    # ── Slide 12: Thank You ──
    s = add_slide(prs)
    set_bg(s, *L_BG)
    rect(s, 0, 0, W, PInches(0.14), L_NAVY)
    rect(s, 0, PInches(2.82), W, PInches(0.06), L_BLUE)
    rect(s, 0, PInches(7.1), W, PInches(0.4), L_NAVY)
    txbox(s, "Thank You", PInches(1), PInches(0.9), PInches(11.3), PInches(1.3),
          size=52, bold=True, color=L_NAVY, align=PP_ALIGN.CENTER)
    txbox(s, "Questions & Discussion",
          PInches(1), PInches(2.92), PInches(11.3), PInches(0.75),
          size=22, color=L_BLUE, align=PP_ALIGN.CENTER, italic=True)
    txbox(s, "MOSIP Nexus  |  AI-Powered Knowledge Assistant  |  May 2026",
          PInches(1), PInches(4.5), PInches(11.3), PInches(0.55),
          size=14, color=L_GRAY, align=PP_ALIGN.CENTER)
    txbox(s, "Built on: LangChain  .  Groq Llama 3.3  .  ChromaDB  .  HuggingFace  .  Streamlit",
          PInches(1), PInches(5.45), PInches(11.3), PInches(0.5),
          size=12, color=L_GRAY, align=PP_ALIGN.CENTER)
    txbox(s, "MOSIP Nexus  |  Confidential",
          PInches(1), PInches(7.13), PInches(11.3), PInches(0.28),
          size=10, color=L_WHITE, align=PP_ALIGN.CENTER)

    out = DOCS_DIR / "MOSIP_Nexus_Presentation.pptx"
    prs.save(str(out))
    print(f"Presentation saved: {out}")


if __name__ == "__main__":
    print("Generating MOSIP Nexus documentation...")
    build_developer_guide()
    build_ppt()
    print("Done. Files are in MosipNexus/docs/")
